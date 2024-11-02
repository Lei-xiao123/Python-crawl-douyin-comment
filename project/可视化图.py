# 设置matplotlib后端
import matplotlib
matplotlib.use('Agg')  # 使用非交互式后端

# 导入数据处理
import pandas as pd
import numpy as np
from PIL import Image
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches

# 设置中文字体
plt.rcParams['font.sans-serif'] = ['SimHei']  # 使用黑体
plt.rcParams['axes.unicode_minus'] = False  # 解决负号显示问题

# 读取文件
df = pd.read_csv('data.csv')

# 按地区分组并统计评论数量
region_counts = df['地区'].value_counts()

# 绘制数据
plt.figure(figsize=(10, 6))
region_counts.plot(kind='bar', color='skyblue')
plt.title('每个地区的评论数量')
plt.xlabel('地区')
plt.ylabel('评论数量')
plt.xticks(rotation=45)
plt.tight_layout()

# 保存图表为图片
plt.savefig('region_comments.png')

# 创建PPT
prs = Presentation()

# 添加幻灯片
slide_layout = prs.slide_layouts[5]  # 使用空白布局
slide = prs.slides.add_slide(slide_layout)

# 添加图片到幻灯片
left = Inches(1)
top = Inches(1)
pic = slide.shapes.add_picture('region_comments.png', left, top, width=Inches(8), height=Inches(4.5))

# 保存PPT
prs.save('comments_by_region.pptx')
